from tkinter import *
from tkinter import filedialog
from tkinter import messagebox
import PyPDF2
import re
import os

from pptx import Presentation


def browse_button():
    global folder_path
    filename = filedialog.askdirectory()
    folder_path.set(filename)


def replace_button():
    pos = 2
    search_path = folder_path.get()
    search_str = str(search.get())

    if not (search_path.endswith("/") or search_path.endswith("\\")):
        search_path = search_path + "/"

    if not folder_path.get():
        messagebox.showerror("Error", "Invalid path")

    for fname in os.listdir(path=search_path):
        if fname.endswith(".ppt") or fname.endswith(".pptx"):
            prs = Presentation(search_path + fname)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape.text = shape.text.lower()
                        if search_str in shape.text:
                            pos += 1
                            list1.insert(pos, '   ' + fname)
                            break
        elif fname.endswith(".pdf"):
            object = PyPDF2.PdfFileReader(search_path + fname)

            # get number of pages
            NumPages = object.getNumPages()

            # define keyterms
            String = search_str

            # extract text and do the search
            for i in range(0, NumPages):
                PageObj = object.getPage(i)
                Text = PageObj.extractText()
                ResSearch = re.search(String, Text)
                if ResSearch:
                    pos += 1
                    list1.insert(pos, '   ' + fname)


root = Tk()
root.configure(background="#3C3F41")
root.title('H file fix')
folder_path = StringVar()

lbl1 = Label(master=root, text='Current Path: ', bg="#3C3F41", fg="#BBBBBB", font='Consolas 11')
lbl1.grid(row=1, column=0)
lbl2 = Label(master=root, textvariable=folder_path, bg="#3C3F41", fg="#BBBBBB", font='Consolas 11')
lbl2.grid(row=1, column=1)

button1 = Button(text="Browse", command=browse_button, width=14, bg="#3C3F41", fg="#BBBBBB", font='Consolas 11')
button1.grid(row=2, column=0, sticky='NW')
button2 = Button(text="Search files", command=replace_button, width=14, bg="#3C3F41", fg="#BBBBBB", font='Consolas 11')
button2.grid(row=2, column=0, sticky='SW')

list1 = Listbox(master=root, height=14, width=72, bg='#2B2B2B', fg="#BBBBBB", font='Consolas 11')
list1.grid(row=2, column=1, rowspan=7, columnspan=8)
list1.insert(0, 'Matching files:')

search = Entry(root, {}, bg='#2B2B2B', fg="#BBBBBB", font='Consolas 11', width=86)
search.grid(row=0, column=0, sticky='N', columnspan=2)

mainloop()
